#!/usr/bin/env python3
"""Generate first-slide cover PNGs for each .pptx in assets/decks/.

Strategy:
1. Try LibreOffice headless: render slide 1 to a PNG. This produces the
   real visual slide and is the preferred path.
2. If LibreOffice fails for a given deck, fall back to a python-pptx
   reconstruction:
   - largest embedded image on slide 1, cover-fit + text overlay
   - or a clean title-card with the actual slide-1 title

Output: assets/decks/thumbs/<original-stem>.png
"""

import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent.parent
DECKS_DIR = ROOT / "assets" / "decks"
OUT_DIR = DECKS_DIR / "thumbs"
OUT_DIR.mkdir(parents=True, exist_ok=True)

W, H = 1280, 800
EQX_RED = (229, 32, 46)
INK = (19, 19, 26)
MUTED = (90, 99, 120)
BG = (246, 245, 240)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    for p in candidates:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def libreoffice_render(pptx_path: Path) -> Image.Image | None:
    """Convert slide 1 to PNG with LibreOffice headless. Return PIL image or None."""
    if not shutil.which("libreoffice"):
        return None
    with tempfile.TemporaryDirectory() as tmpdir:
        # LO needs a writable HOME for its profile
        env = {**os.environ, "HOME": tmpdir}
        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--norestore",
                    "--nologo",
                    "--convert-to", "png",
                    "--outdir", tmpdir,
                    str(pptx_path),
                ],
                env=env,
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                print(f"     [LO failed] {pptx_path.name}: {result.stderr.strip()[:120]}", file=sys.stderr)
                return None
        except Exception as e:
            print(f"     [LO error] {pptx_path.name}: {e}", file=sys.stderr)
            return None

        out = Path(tmpdir) / (pptx_path.stem + ".png")
        if not out.exists():
            return None
        return Image.open(out).convert("RGB").copy()


def cover_fit(img: Image.Image, w: int, h: int) -> Image.Image:
    img = img.convert("RGB")
    iw, ih = img.size
    ir = iw / ih
    cr = w / h
    if ir > cr:
        new_h = h
        new_w = int(h * ir)
    else:
        new_w = w
        new_h = int(w / ir)
    img = img.resize((new_w, new_h), Image.LANCZOS)
    dx = (new_w - w) // 2
    dy = (new_h - h) // 2
    return img.crop((dx, dy, dx + w, dy + h))


def upscale_to_canvas(img: Image.Image) -> Image.Image:
    """Place the rendered slide on a 1280x800 canvas. Cover-fit and crop centered."""
    return cover_fit(img, W, H)


def wrap_text(draw, text: str, font, max_w: int) -> list[str]:
    words = text.split()
    lines, cur = [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if draw.textlength(trial, font=font) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip()
        if t:
            return t
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().splitlines()[0]
    return ""


def largest_image(slide):
    best = None
    best_area = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            try:
                blob = shape.image.blob
                img = Image.open(io.BytesIO(blob))
                area = img.width * img.height
                if area > best_area:
                    best_area, best = area, img
            except Exception:
                continue
    return best


def render_image_card(title: str, author: str, img: Image.Image) -> Image.Image:
    canvas = cover_fit(img, W, H).copy()
    draw = ImageDraw.Draw(canvas, "RGBA")
    grad = Image.new("L", (1, H))
    for y in range(H):
        if y < H - 260:
            grad.putpixel((0, y), 0)
        else:
            t = (y - (H - 260)) / 260
            grad.putpixel((0, y), int(220 * (t ** 1.4)))
    grad = grad.resize((W, H))
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    overlay.putalpha(grad)
    canvas.paste(Image.new("RGB", (W, H), (0, 0, 0)), (0, 0), overlay)
    draw.rectangle((0, 0, 14, H), fill=EQX_RED)
    f_title = load_font(56, bold=True)
    f_author = load_font(28)
    f_eyebrow = load_font(22, bold=True)
    draw.text((44, H - 230), "STRATEGY DECK · GTM AI", font=f_eyebrow, fill=(255, 255, 255, 230))
    lines = wrap_text(draw, title, f_title, W - 100)[:2]
    y = H - 180
    for line in lines:
        draw.text((44, y), line, font=f_title, fill="white")
        y += 64
    if author:
        draw.text((44, H - 56), f"by {author}", font=f_author, fill=(220, 224, 235))
    return canvas


def render_text_card(title: str, author: str) -> Image.Image:
    canvas = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(canvas)
    for x in range(0, W, 64):
        draw.line([(x, 0), (x, H)], fill=(0, 0, 0, 8), width=1)
    draw.rectangle((0, 0, 14, H), fill=EQX_RED)
    f_eyebrow = load_font(22, bold=True)
    draw.text((56, 100), "STRATEGY DECK · GTM AI", font=f_eyebrow, fill=EQX_RED)
    f_title = load_font(70, bold=True)
    lines = wrap_text(draw, title, f_title, W - 110)[:3]
    y = 200
    for line in lines:
        draw.text((56, y), line, font=f_title, fill=INK)
        y += 80
    if author:
        f_author = load_font(28)
        draw.text((56, H - 90), f"by {author}", font=f_author, fill=MUTED)
    f_brand = load_font(20, bold=True)
    f_meta = load_font(16)
    brand = "EQUINIX · GTM ENABLEMENT"
    bw = draw.textlength(brand, font=f_brand)
    draw.text((W - 56 - bw, H - 90), brand, font=f_brand, fill=INK)
    sw = draw.textlength("Slide 1", font=f_meta)
    draw.text((W - 56 - sw, H - 56), "Slide 1", font=f_meta, fill=(154, 161, 173))
    return canvas


AUTHOR_MAP = {
    "AI CRO Strategy Plan - Rizwan Qureshy.pptx": "Rizwan Qureshy",
    "AI Strategy - Eamonn Deck.pptx": "Eamonn Ward",
    "AI Strategy - Morgan Gallegos.pptx": "Morgan Gallegos",
    "AI Strategy Deck - Justin Sit.pptx": "Justin Sit",
    "AI Strategy Deck_April 3rd - Paige Gregory.pptx": "Paige Gregory",
    "AI for GTM Enablement Services - Kelly Grover.pptx": "Kelly Grover",
    "AI in our team - Calley Hood.pptx": "Calley Hood",
    "AI-Strategy-for-GTM-Enablement - Isabelle Puller.pptx": "Isabelle Puller",
    "AI-Strategy-for-Scalable-Partner-Enablement - Michael Bourgeois.pptx": "Michael Bourgeois",
    "AI_Strategy_Deck_claude - Jason Sherwood.pptx": "Jason Sherwood",
    "AI_in_Our_Team_Strategy - Lorna Joiner.pptx": "Lorna Joiner",
}


def main():
    decks = sorted(DECKS_DIR.glob("*.pptx"))
    print(f"Found {len(decks)} decks")
    for deck in decks:
        name = deck.name
        out_path = OUT_DIR / (deck.stem + ".png")

        # 1. Try LibreOffice
        rendered = libreoffice_render(deck)
        if rendered is not None:
            card = upscale_to_canvas(rendered)
            card.save(out_path, "PNG", optimize=True)
            print(f"  ok  LO     {name}")
            continue

        # 2. Fall back to python-pptx-based reconstruction
        try:
            pres = Presentation(str(deck))
            slide = pres.slides[0]
            title = re.sub(r"\s+", " ", slide_title(slide) or deck.stem).strip()
            author = AUTHOR_MAP.get(name, "Team")
            img = largest_image(slide)
            if img is not None and img.width >= 320 and img.height >= 240:
                card = render_image_card(title, author, img)
                kind = "img-fb"
            else:
                card = render_text_card(title, author)
                kind = "txt-fb"
            card.save(out_path, "PNG", optimize=True)
            print(f"  ok  {kind}  {name}")
        except Exception as e:
            print(f"  FAIL          {name}: {e}", file=sys.stderr)


if __name__ == "__main__":
    main()

